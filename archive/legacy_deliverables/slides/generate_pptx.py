#!/usr/bin/env python3
"""Generate the internship review presentation as a PowerPoint file."""

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# Create presentation
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Colors
DARK_BLUE = RGBColor(0x1A, 0x23, 0x7E)
ACCENT = RGBColor(0x00, 0x96, 0xC7)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
LIGHT_GRAY = RGBColor(0x66, 0x66, 0x66)
GREEN = RGBColor(0x28, 0xA7, 0x45)
ORANGE = RGBColor(0xF5, 0x7C, 0x00)


def add_title_slide(prs, title, subtitle):
    slide_layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(slide_layout)

    # Background shape
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK_BLUE
    shape.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.3), Inches(1.2))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.8), Inches(12.3), Inches(1))
    tf = sub_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(24)
    p.font.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
    p.alignment = PP_ALIGN.CENTER

    # Footer
    foot_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.8), Inches(12.3), Inches(0.4))
    tf = foot_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Srujan | SWA Consultancy Pvt. Ltd. | June 2026"
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(0x99, 0x99, 0x99)
    p.alignment = PP_ALIGN.CENTER

    return slide


def add_content_slide(prs, title, bullets, notes=""):
    slide_layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(slide_layout)

    # Header bar
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.1))
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK_BLUE
    shape.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(12.3), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # Content
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(12.3), Inches(5.5))
    tf = content_box.text_frame
    tf.word_wrap = True

    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = bullet
        p.font.size = Pt(20)
        p.font.color.rgb = DARK_GRAY
        p.space_before = Pt(12)
        p.space_after = Pt(6)
        p.level = 0

    # Notes
    if notes:
        slide.notes_slide.notes_text_frame.text = notes

    return slide


def add_table_slide(prs, title, headers, rows, notes=""):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Header bar
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.1))
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK_BLUE
    shape.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(12.3), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # Table
    num_rows = len(rows) + 1
    num_cols = len(headers)
    table = slide.shapes.add_table(
        num_rows, num_cols, Inches(0.5), Inches(1.4), Inches(12.3), Inches(0.6 * num_rows)
    ).table

    # Header row
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = ACCENT
        paragraph = cell.text_frame.paragraphs[0]
        paragraph.font.size = Pt(14)
        paragraph.font.bold = True
        paragraph.font.color.rgb = WHITE

    # Data rows
    for row_idx, row_data in enumerate(rows):
        for col_idx, cell_text in enumerate(row_data):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = str(cell_text)
            paragraph = cell.text_frame.paragraphs[0]
            paragraph.font.size = Pt(13)
            paragraph.font.color.rgb = DARK_GRAY
            if col_idx == len(row_data) - 1 and isinstance(cell_text, str):
                if "Strong" in cell_text or "Done" in cell_text:
                    paragraph.font.color.rgb = GREEN
                elif "Bug" in cell_text or "Hang" in cell_text:
                    paragraph.font.color.rgb = RGBColor(0xDC, 0x35, 0x45)
                elif "Fix" in cell_text or "Tuning" in cell_text:
                    paragraph.font.color.rgb = ORANGE

    if notes:
        slide.notes_slide.notes_text_frame.text = notes

    return slide


def add_two_column_slide(prs, title, left_title, left_bullets, right_title, right_bullets, notes=""):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Header bar
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.1))
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK_BLUE
    shape.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(12.3), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # Left column title
    lt_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(5.9), Inches(0.5))
    tf = lt_box.text_frame
    p = tf.paragraphs[0]
    p.text = left_title
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = ACCENT

    # Left bullets
    lb_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.9), Inches(5.9), Inches(5))
    tf = lb_box.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(left_bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = bullet
        p.font.size = Pt(16)
        p.font.color.rgb = DARK_GRAY
        p.space_before = Pt(8)

    # Right column title
    rt_box = slide.shapes.add_textbox(Inches(6.9), Inches(1.4), Inches(5.9), Inches(0.5))
    tf = rt_box.text_frame
    p = tf.paragraphs[0]
    p.text = right_title
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = ACCENT

    # Right bullets
    rb_box = slide.shapes.add_textbox(Inches(6.9), Inches(1.9), Inches(5.9), Inches(5))
    tf = rb_box.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(right_bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = bullet
        p.font.size = Pt(16)
        p.font.color.rgb = DARK_GRAY
        p.space_before = Pt(8)

    if notes:
        slide.notes_slide.notes_text_frame.text = notes

    return slide


# === SLIDE 1: Title ===
add_title_slide(prs, "RFQ2BOQ", "NLP-Based Extraction of Bill of Quantities\nfrom Construction Tenders")

# === SLIDE 2: Introduction ===
add_content_slide(
    prs,
    "Introduction — The Problem",
    [
        "Construction tenders (RFQs) contain BOQ data buried in PDFs and Excel files",
        "Manual extraction by estimators takes 2–4 hours per tender",
        "Process is slow, error-prone, and doesn't scale",
        "",
        "Goal: Build an AI system that reads tenders and outputs structured BOQ automatically",
        "Extract 8 entity types: Material, Quantity, Unit, Location, Dimension, Standard, Grade, Action",
        "Export as Excel (CPWD format), JSON, or CSV",
    ],
    "Keep it relatable. Everyone understands reading PDFs is painful.",
)

# === SLIDE 3: Objectives ===
add_two_column_slide(
    prs,
    "Internship Objectives",
    "Primary",
    [
        "End-to-end NLP pipeline for BOQ extraction",
        "Handle both PDF and Excel tender formats",
        "Extract 8 entity types with relations",
        "Export structured BOQ (Excel/JSON/CSV)",
    ],
    "Secondary",
    [
        "Streamlit UI for non-technical users",
        "FastAPI backend for integration",
        "Honest evaluation framework (no fake 100%)",
        "Full documentation for future interns",
    ],
    "Show clear scope from day one.",
)

# === SLIDE 4: Architecture ===
add_content_slide(
    prs,
    "System Architecture",
    [
        "PDF / XLSX  →  Ingest  →  Preprocess  →  NER  →  Domain  →  Export",
        "",
        "Ingest: pdfplumber for tables + pytesseract OCR for scanned PDFs",
        "Preprocess: Text cleaning + SmartSectionClassifier (finds BOQ pages)",
        "NER: Pattern-based (production) + BERT-LoRA ML (experimental)",
        "Domain: BOQ assembly + unit normalization + validation",
        "Export: Excel (CPWD), JSON, CSV",
        "",
        "Tech Stack: Python, PyTorch, Transformers, pdfplumber, FastAPI, Streamlit",
    ],
    "Walk left to right. Don't rush the diagram.",
)

# === SLIDE 5: What Was Built ===
add_table_slide(
    prs,
    "What Was Built — Components Delivered",
    ["Component", "Status", "Details"],
    [
        ["PDF Ingestion", "Done", "pdfplumber + OCR fallback"],
        ["XLSX Ingestion", "Done", "Row-preservation pipeline"],
        ["NER Engine", "Done", "Pattern-based (prod) + LoRA ML (exp)"],
        ["BOQ Assembler", "Done", "Row build + unit normalization"],
        ["Validation", "Done", "Domain rules + confidence scoring"],
        ["Export", "Done", "Excel, JSON, CSV (CPWD format)"],
        ["CLI / API / UI", "Done", "Typer, FastAPI, Streamlit"],
        ["Tests", "Done", "97 critical tests passing"],
        ["Evaluation", "Done", "Honest, independent gold"],
        ["Documentation", "Done", "Handoff, architecture, user guide"],
    ],
    "Emphasize dual NER: pattern-based for reliability today, ML for future.",
)

# === SLIDE 6: Evaluation ===
add_content_slide(
    prs,
    "Honest Evaluation Framework",
    [
        "Entity-level F1: 43.8% micro / 44.8% macro (independent gold)",
        "XLSX extraction: Strong — 60% macro F1, exact row counts on clean spreadsheets",
        "PDF extraction: Partial — 23% macro F1, needs improvement",
        "",
        "Important Caveat — Evaluation Mismatch:",
        "  • Gold expects short material names: 'Mineral Wool'",
        "  • Pipeline outputs full descriptions: 'Supply & application of 100 mm thick Mineral Wool...'",
        "  • This makes F1 look worse than actual row correctness",
        "",
        "Anti-Cheat Rules Implemented:",
        "  • No self-comparison (pipeline output ≠ gold)",
        "  • Independent human-verified rowgold",
        "  • Automatic detection of gold modification",
    ],
    "Be honest about numbers but explain WHY. Don't hide the mismatch.",
)

# === SLIDE 7: Demo Table ===
add_table_slide(
    prs,
    "Demo — 10 Real SWA Tender Files",
    ["File", "Type", "Items", "Status"],
    [
        ["02 ISRO", "XLSX", "8 rows", "Strong"],
        ["03 Zydus Matoda", "XLSX", "33 rows", "Strong"],
        ["05 Zydus Animal", "XLSX", "48 rows", "Strong"],
        ["08 SAEL", "XLSX", "12 rows", "Strong"],
        ["04 Adani", "PDF", "2 (bug)", "Fixing"],
        ["06 Avante", "PDF", "31 (some FP)", "Tuning"],
        ["07 Grew", "PDF", "9 (some FP)", "Tuning"],
        ["01 GSECL", "PDF", "2 (weak)", "Fixing"],
        ["09 GeM", "PDF", "22 (slow)", "Optimizing"],
        ["10 GeM", "PDF", "10", "OK"],
    ],
    "Lead with wins (XLSX), then be honest about PDF gaps.",
)

# === SLIDE 8: Learnings ===
add_two_column_slide(
    prs,
    "Key Learnings",
    "Technical",
    [
        "Data quality beats architecture",
        "Synthetic 99% F1 → real 43% F1",
        "Honest evaluation is non-negotiable",
        "PDF is harder than Excel",
        "Pattern NER beats ML on small data",
    ],
    "Process",
    [
        "One agent at a time on repo",
        "Every fix committed before next",
        "Say no to out-of-scope features",
        "Scope guard prevents drift",
    ],
    "Show maturity — you learned from mistakes.",
)

# === SLIDE 9: Challenges ===
add_content_slide(
    prs,
    "Challenges Faced",
    [
        "1. Synthetic vs Real Data Gap",
        "   • Model trained on 300 synthetic PDFs → 99% F1",
        "   • Same model on real tenders → MATERIAL F1 = 0.0",
        "   • Root cause: synthetic data was regex-generated from research papers",
        "",
        "2. PDF Table Fragility",
        "   • Merged cells, multi-line cells, split-quantity columns",
        "   • GeM PDFs have digit columns text classifiers miss",
        "",
        "3. The 'Fake 100%' Trap",
        "   • Earlier work modified gold files to match pipeline output",
        "   • We caught this and implemented anti-cheat rules",
        "   • All numbers now independently verified",
    ],
    "Don't blame others. Say 'we caught and fixed it.'",
)

# === SLIDE 10: Tools ===
add_two_column_slide(
    prs,
    "Tools & Skills Developed",
    "Technologies",
    [
        "Python 3.11–3.13",
        "PyTorch + HuggingFace Transformers",
        "pdfplumber + pytesseract (OCR)",
        "openpyxl (Excel)",
        "FastAPI + Streamlit",
        "pytest + ruff + Docker",
    ],
    "Skills",
    [
        "End-to-end ML pipeline design",
        "Domain-specific NLP",
        "PDF extraction techniques",
        "Honest evaluation methodology",
        "Multi-agent project management",
        "Git discipline & code review",
    ],
    "Show breadth — full stack understanding.",
)

# === SLIDE 11: Achievements ===
add_content_slide(
    prs,
    "Achievements & Outcomes",
    [
        "Delivered: Complete pipeline — PDF/XLSX → structured BOQ in <30 seconds",
        "Delivered: Working UI + API + CLI for multiple user types",
        "Delivered: Honest evaluation framework with anti-cheat safeguards",
        "Delivered: 97 critical tests passing with CI/CD gate",
        "Delivered: Full documentation (architecture, handoff, user guide)",
        "",
        "Honest Metrics:",
        "  • XLSX: Strong — exact row counts on 4/4 files",
        "  • Entity F1: 43.8% micro / 44.8% macro",
        "  • PDF: Partial — known bugs being fixed",
        "  • Crash rate: 0% — all 10 files process end-to-end",
        "",
        "Key Takeaway: Foundation is solid. Bottleneck is data, not code.",
    ],
    "End on positive but honest note.",
)

# === SLIDE 12: Conclusion ===
add_content_slide(
    prs,
    "Conclusion & Future Work",
    [
        "Summary: Built complete RFQ→BOQ system. XLSX works today. PDF needs real training data.",
        "",
        "Immediate Next Steps:",
        "  1. Fix PDF bugs (Adani headers, GSECL page detection)",
        "  2. Human annotation: 20–40 real tenders",
        "  3. NER retrain on real gold → target F1 > 0.60",
        "",
        "Long-Term Vision:",
        "  • Domain-specific models (insulation, civil, electrical)",
        "  • Hindi/Indic language support (module ready)",
        "  • Batch processing for high-volume use",
        "  • Human-in-the-loop for low-confidence extractions",
        "",
        "Thank You — Questions?",
    ],
    "Clear summary + confident close.",
)

# Save
output_path = "/Users/srujansai/Desktop/rfq2boq/deliverables/slides/Internship_Review_Presentation.pptx"
prs.save(output_path)
print(f"Presentation saved to: {output_path}")
print(f"Total slides: {len(prs.slides)}")
